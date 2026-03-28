import fitz
import os
from docx import Document
from pptx import Presentation

from dependencies import get_embedding_model


def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract text from PDF using PyMuPDF"""
    try:
        doc = fitz.open(pdf_path)
        text = ""
        for page in doc:
            text += page.get_text()
        return text
    except Exception as e:
        print(f"PDF extraction error: {e}")
        return ""


def extract_text_from_docx(docx_path: str) -> str:
    """Extract text from DOCX file"""
    try:
        doc = Document(docx_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        return text
    except Exception as e:
        print(f"DOCX extraction error: {e}")
        return ""


def extract_text_from_pptx(pptx_path: str) -> str:
    """Extract text from PPTX file"""
    try:
        prs = Presentation(pptx_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"PPTX extraction error: {e}")
        return ""


def extract_text_from_txt(txt_path: str) -> str:
    """Extract text from TXT file"""
    try:
        with open(txt_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            with open(txt_path, "r", encoding="latin-1") as f:
                return f.read()
        except Exception as e:
            print(f"TXT extraction error (latin-1): {e}")
            return ""
    except Exception as e:
        print(f"TXT extraction error: {e}")
        return ""


def extract_text_by_file_type(file_path: str, file_extension: str) -> str:
    """Route to appropriate extraction function based on file type"""
    ext = file_extension.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    elif ext == ".txt":
        return extract_text_from_txt(file_path)
    else:
        print(f"Unsupported file type: {ext}")
        return ""


def chunk_text(text: str, chunk_size: int = 1500, overlap: int = 200) -> list[str]:
    """Split text into overlapping chunks"""
    chunks: list[str] = []
    start = 0

    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end - overlap

    return chunks


def embed_text(text: str) -> list[float]:
    try:
        embedding_model = get_embedding_model()
        embedding = embedding_model.encode(text).tolist()
        return embedding
    except Exception as e:
        print(f"Embedding error: {e}")
        return []
